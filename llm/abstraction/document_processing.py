"""
Copyright 2025-2030 all rights reserved
Ashutosh Sinha
Email: ajsinha@gmail.com
Version: 3.1.4

Document Processing - PDF, DOCX, XLSX, PPT, Markdown, HTML, etc.
"""

import os
from typing import Optional, Dict, Any, List, Union
from dataclasses import dataclass
from pathlib import Path
from enum import Enum


class DocumentType(Enum):
    """Supported document types"""
    PDF = "pdf"
    DOCX = "docx"
    DOC = "doc"
    XLSX = "xlsx"
    XLS = "xls"
    PPTX = "pptx"
    PPT = "ppt"
    TXT = "txt"
    MD = "markdown"
    HTML = "html"
    CSV = "csv"
    JSON = "json"
    XML = "xml"
    RTF = "rtf"


@dataclass
class DocumentContent:
    """
    Represents processed document content.
    
    Attributes:
        text: Extracted text content
        doc_type: Type of document
        metadata: Document metadata (title, author, pages, etc.)
        structure: Optional structured data (tables, sections, etc.)
        images: Optional extracted images
    """
    text: str
    doc_type: DocumentType
    metadata: Dict[str, Any]
    structure: Optional[Dict[str, Any]] = None
    images: Optional[List[bytes]] = None
    
    def __post_init__(self):
        if self.structure is None:
            self.structure = {}
        if self.images is None:
            self.images = []


class DocumentProcessor:
    """
    Universal document processor for extracting content from various formats.
    
    Example:
        processor = DocumentProcessor()
        
        # Process PDF
        content = processor.process("document.pdf")
        print(content.text)
        print(content.metadata)
        
        # Process with options
        content = processor.process(
            "report.docx",
            extract_images=True,
            extract_tables=True
        )
    """
    
    SUPPORTED_FORMATS = {
        '.pdf': DocumentType.PDF,
        '.docx': DocumentType.DOCX,
        '.doc': DocumentType.DOC,
        '.xlsx': DocumentType.XLSX,
        '.xls': DocumentType.XLS,
        '.pptx': DocumentType.PPTX,
        '.ppt': DocumentType.PPT,
        '.txt': DocumentType.TXT,
        '.md': DocumentType.MD,
        '.html': DocumentType.HTML,
        '.htm': DocumentType.HTML,
        '.csv': DocumentType.CSV,
        '.json': DocumentType.JSON,
        '.xml': DocumentType.XML,
        '.rtf': DocumentType.RTF,
    }
    
    def __init__(self):
        self.processors = {
            DocumentType.PDF: self._process_pdf,
            DocumentType.DOCX: self._process_docx,
            DocumentType.XLSX: self._process_xlsx,
            DocumentType.PPTX: self._process_pptx,
            DocumentType.TXT: self._process_txt,
            DocumentType.MD: self._process_markdown,
            DocumentType.HTML: self._process_html,
            DocumentType.CSV: self._process_csv,
            DocumentType.JSON: self._process_json,
        }
    
    def process(
        self,
        filepath: str,
        extract_images: bool = False,
        extract_tables: bool = False,
        page_range: Optional[tuple] = None
    ) -> DocumentContent:
        """
        Process a document and extract content.
        
        Args:
            filepath: Path to document
            extract_images: Extract images from document
            extract_tables: Extract tables as structured data
            page_range: Optional (start, end) page range
        
        Returns:
            DocumentContent object
        """
        if not os.path.exists(filepath):
            raise FileNotFoundError(f"Document not found: {filepath}")
        
        ext = Path(filepath).suffix.lower()
        doc_type = self.SUPPORTED_FORMATS.get(ext)
        
        if not doc_type:
            raise ValueError(f"Unsupported document format: {ext}")
        
        processor = self.processors.get(doc_type)
        if not processor:
            raise NotImplementedError(f"Processor for {doc_type} not implemented")
        
        return processor(
            filepath,
            extract_images=extract_images,
            extract_tables=extract_tables,
            page_range=page_range
        )
    
    def _process_pdf(
        self,
        filepath: str,
        extract_images: bool = False,
        extract_tables: bool = False,
        page_range: Optional[tuple] = None
    ) -> DocumentContent:
        """Process PDF document"""
        try:
            import PyPDF2
            
            text = []
            metadata = {}
            images = []
            
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                
                # Extract metadata
                if reader.metadata:
                    metadata = {
                        'title': reader.metadata.get('/Title', ''),
                        'author': reader.metadata.get('/Author', ''),
                        'pages': len(reader.pages)
                    }
                
                # Extract text
                start = page_range[0] if page_range else 0
                end = page_range[1] if page_range else len(reader.pages)
                
                for i in range(start, min(end, len(reader.pages))):
                    page = reader.pages[i]
                    text.append(page.extract_text())
            
            return DocumentContent(
                text='\n\n'.join(text),
                doc_type=DocumentType.PDF,
                metadata=metadata,
                images=images if extract_images else None
            )
        
        except ImportError:
            # Fallback: Simple text extraction
            return DocumentContent(
                text=f"[PDF content from {filepath}]\nInstall PyPDF2 for full support: pip install PyPDF2",
                doc_type=DocumentType.PDF,
                metadata={'filename': os.path.basename(filepath)}
            )
    
    def _process_docx(
        self,
        filepath: str,
        extract_images: bool = False,
        extract_tables: bool = False,
        **kwargs
    ) -> DocumentContent:
        """Process DOCX document"""
        try:
            import docx
            
            doc = docx.Document(filepath)
            
            # Extract text
            text = []
            for para in doc.paragraphs:
                text.append(para.text)
            
            # Extract metadata
            metadata = {
                'filename': os.path.basename(filepath),
                'paragraphs': len(doc.paragraphs)
            }
            
            # Extract tables
            structure = {}
            if extract_tables and doc.tables:
                tables = []
                for table in doc.tables:
                    table_data = []
                    for row in table.rows:
                        table_data.append([cell.text for cell in row.cells])
                    tables.append(table_data)
                structure['tables'] = tables
            
            return DocumentContent(
                text='\n'.join(text),
                doc_type=DocumentType.DOCX,
                metadata=metadata,
                structure=structure
            )
        
        except ImportError:
            return DocumentContent(
                text=f"[DOCX content from {filepath}]\nInstall python-docx: pip install python-docx",
                doc_type=DocumentType.DOCX,
                metadata={'filename': os.path.basename(filepath)}
            )
    
    def _process_xlsx(
        self,
        filepath: str,
        extract_tables: bool = True,
        **kwargs
    ) -> DocumentContent:
        """Process XLSX spreadsheet"""
        try:
            import openpyxl
            
            wb = openpyxl.load_workbook(filepath)
            
            text = []
            structure = {'sheets': []}
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text.append(f"\n=== {sheet_name} ===\n")
                
                sheet_data = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = '\t'.join(str(cell) if cell else '' for cell in row)
                        text.append(row_text)
                        sheet_data.append(list(row))
                
                if extract_tables:
                    structure['sheets'].append({
                        'name': sheet_name,
                        'data': sheet_data
                    })
            
            return DocumentContent(
                text='\n'.join(text),
                doc_type=DocumentType.XLSX,
                metadata={
                    'filename': os.path.basename(filepath),
                    'sheets': len(wb.sheetnames)
                },
                structure=structure
            )
        
        except ImportError:
            return DocumentContent(
                text=f"[XLSX content from {filepath}]\nInstall openpyxl: pip install openpyxl",
                doc_type=DocumentType.XLSX,
                metadata={'filename': os.path.basename(filepath)}
            )
    
    def _process_pptx(
        self,
        filepath: str,
        extract_images: bool = False,
        **kwargs
    ) -> DocumentContent:
        """Process PPTX presentation"""
        try:
            from pptx import Presentation
            
            prs = Presentation(filepath)
            
            text = []
            structure = {'slides': []}
            
            for i, slide in enumerate(prs.slides, 1):
                text.append(f"\n=== Slide {i} ===\n")
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                text.extend(slide_text)
                structure['slides'].append({
                    'number': i,
                    'content': slide_text
                })
            
            return DocumentContent(
                text='\n'.join(text),
                doc_type=DocumentType.PPTX,
                metadata={
                    'filename': os.path.basename(filepath),
                    'slides': len(prs.slides)
                },
                structure=structure
            )
        
        except ImportError:
            return DocumentContent(
                text=f"[PPTX content from {filepath}]\nInstall python-pptx: pip install python-pptx",
                doc_type=DocumentType.PPTX,
                metadata={'filename': os.path.basename(filepath)}
            )
    
    def _process_txt(self, filepath: str, **kwargs) -> DocumentContent:
        """Process plain text file"""
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        
        return DocumentContent(
            text=text,
            doc_type=DocumentType.TXT,
            metadata={'filename': os.path.basename(filepath)}
        )
    
    def _process_markdown(self, filepath: str, **kwargs) -> DocumentContent:
        """Process Markdown file"""
        with open(filepath, 'r', encoding='utf-8') as f:
            text = f.read()
        
        return DocumentContent(
            text=text,
            doc_type=DocumentType.MD,
            metadata={'filename': os.path.basename(filepath)}
        )
    
    def _process_html(self, filepath: str, **kwargs) -> DocumentContent:
        """Process HTML file"""
        try:
            from bs4 import BeautifulSoup
            
            with open(filepath, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
            
            # Extract text
            text = soup.get_text(separator='\n', strip=True)
            
            # Extract metadata
            metadata = {'filename': os.path.basename(filepath)}
            if soup.title:
                metadata['title'] = soup.title.string
            
            return DocumentContent(
                text=text,
                doc_type=DocumentType.HTML,
                metadata=metadata
            )
        
        except ImportError:
            with open(filepath, 'r', encoding='utf-8') as f:
                text = f.read()
            
            return DocumentContent(
                text=text,
                doc_type=DocumentType.HTML,
                metadata={'filename': os.path.basename(filepath)}
            )
    
    def _process_csv(self, filepath: str, **kwargs) -> DocumentContent:
        """Process CSV file"""
        import csv
        
        text = []
        structure = {'rows': []}
        
        with open(filepath, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                text.append('\t'.join(row))
                structure['rows'].append(row)
        
        return DocumentContent(
            text='\n'.join(text),
            doc_type=DocumentType.CSV,
            metadata={'filename': os.path.basename(filepath)},
            structure=structure
        )
    
    def _process_json(self, filepath: str, **kwargs) -> DocumentContent:
        """Process JSON file"""
        import json
        
        with open(filepath, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        return DocumentContent(
            text=json.dumps(data, indent=2),
            doc_type=DocumentType.JSON,
            metadata={'filename': os.path.basename(filepath)},
            structure={'data': data}
        )
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file extensions"""
        return list(self.SUPPORTED_FORMATS.keys())
    
    def is_supported(self, filepath: str) -> bool:
        """Check if file format is supported"""
        ext = Path(filepath).suffix.lower()
        return ext in self.SUPPORTED_FORMATS


def batch_process_documents(
    directory: str,
    pattern: str = "*",
    recursive: bool = True
) -> List[DocumentContent]:
    """
    Batch process multiple documents in a directory.
    
    Args:
        directory: Directory path
        pattern: File pattern (e.g., "*.pdf")
        recursive: Search subdirectories
    
    Returns:
        List of DocumentContent objects
    """
    processor = DocumentProcessor()
    results = []
    
    path = Path(directory)
    files = path.rglob(pattern) if recursive else path.glob(pattern)
    
    for filepath in files:
        if processor.is_supported(str(filepath)):
            try:
                content = processor.process(str(filepath))
                results.append(content)
            except Exception as e:
                print(f"Error processing {filepath}: {e}")
    
    return results


__all__ = [
    'DocumentType',
    'DocumentContent',
    'DocumentProcessor',
    'batch_process_documents',
]


"""
Copyright 2025-2030 all rights reserved
Ashutosh Sinha | ajsinha@gmail.com | Version: 3.1.4
"""
